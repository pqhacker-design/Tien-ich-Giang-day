import io
from pptx import Presentation
from pptx.util import Pt

class PPTXProcessor:
    @staticmethod
    def extract_slides_text(file_bytes: bytes) -> list[dict]:
        """
        Trích xuất nội dung văn bản của từng slide kèm số thứ tự slide.
        Trả về danh sách các slide: [{'slide_number': 1, 'content': '...'}, ...]
        """
        prs = Presentation(io.BytesIO(file_bytes))
        slides_data = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            # Đọc nội dung trên các shapes/khung chữ của slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_texts.append(text)

            content = "\n".join(slide_texts).strip()
            if content:
                slides_data.append({
                    "slide_number": idx,
                    "content": content
                })
        
        return slides_data

    @staticmethod
    def format_doc_text_for_ai(slides_data: list[dict]) -> str:
        """Định dạng dữ liệu các slide thành văn bản có phân tách rõ ràng để AI phân tích."""
        formatted_parts = []
        for item in slides_data:
            formatted_parts.append(
                f"=== SLIDE {item['slide_number']} ===\n{item['content']}\n"
            )
        return "\n".join(formatted_parts)

    @staticmethod
    def integrate_into_notes(file_bytes: bytes, ai_data: dict) -> io.BytesIO:
        """
        Chèn gợi ý Năng lực số / Năng lực AI vào phần Speaker Notes của từng slide tương ứng.
        """
        prs = Presentation(io.BytesIO(file_bytes))
        sua_doi_list = ai_data.get('sua_doi', [])
        
        # Gom nhóm các nội dung tích hợp theo số thứ tự slide
        notes_by_slide = {}
        for item in sua_doi_list:
            slide_num = item.get('slide_number')
            content = item.get('insert_content', '').strip()
            loai = item.get('loai', 'Năng lực số')
            
            if slide_num is not None and content:
                tag = "[NĂNG LỰC AI]" if loai == "Năng lực AI" else "[NĂNG LỰC SỐ]"
                note_entry = f"{tag}: {content}"
                notes_by_slide.setdefault(int(slide_num), []).append(note_entry)

        # Duyệt qua các slide trong presentation và chèn vào Notes
        for idx, slide in enumerate(prs.slides, start=1):
            if idx in notes_by_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                
                existing_text = text_frame.text.strip()
                integration_text = "\n\n".join(notes_by_slide[idx])
                
                header_sep = "📌 [TÍCH HỢP NĂNG LỰC SỐ & AI]:\n"
                
                if existing_text:
                    text_frame.text = f"{existing_text}\n\n{header_sep}{integration_text}"
                else:
                    text_frame.text = f"{header_sep}{integration_text}"

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream
```eof

```python:gemini_service.py
import json
from pydantic import BaseModel, Field
from typing import List, Literal, Optional
from google import genai
from google.genai import types
from google.genai.errors import APIError

# Schema dữ liệu hỗ trợ cả Word (anchor_text) và PowerPoint (slide_number)
class SuaDoiItem(BaseModel):
    slide_number: Optional[int] = Field(default=None, description="Số thứ tự của Slide trong file PowerPoint (nếu là bài giảng PPTX).")
    anchor_text: Optional[str] = Field(default="", description="Câu văn/dòng neo có thật trong giáo án Word.")
    insert_content: str = Field(description="Nội dung sư phạm tích hợp ngắn gọn, bắt đầu bằng hành động của HS hoặc hướng dẫn của GV.")
    loai: Literal["Năng lực số", "Năng lực AI"] = Field(description="Loại năng lực được tích hợp: 'Năng lực số' hoặc 'Năng lực AI'.")

class TichHopResult(BaseModel):
    sua_doi: List[SuaDoiItem]


class GeminiService:
    def __init__(self, api_key: str):
        if not api_key:
            raise ValueError("API Key không được để trống.")
        self.client = genai.Client(api_key=api_key)
        self.model_name = "gemini-2.5-flash"

    def _get_tt02_framework_prompt(self, cap_hoc: str) -> str:
        base_framework = """
* Khung Năng lực số (TT 02/2025/TT-BGDĐT):
  1. Vận hành thiết bị và phần mềm
  2. Khai thác thông tin và dữ liệu
  3. Giao tiếp và hợp tác trong môi trường số
  4. Sáng tạo nội dung số
  5. An toàn và an ninh số
  6. Giải quyết vấn đề trong môi trường số
"""
        level_guide = {
            "Tiểu học": "\n- Tiểu học: Thao tác cơ bản, tìm kiếm đơn giản, ý thức bảo vệ tư thế, mắt và thông tin cá nhân.",
            "THCS": "\n- THCS: Khai thác phần mềm môn học, đánh giá thông tin, làm việc nhóm trực tuyến an toàn, tôn trọng bản quyền.",
            "THPT": "\n- THPT: Xử lý và phân tích dữ liệu, tạo sản phẩm số, an toàn thông tin và ứng dụng công nghệ giải quyết vấn đề.",
            "Tự động nhận diện": "\n- Tự nhận diện cấp học theo từng nội dung để tích hợp vừa sức học sinh."
        }
        return base_framework + level_guide.get(cap_hoc, level_guide["Tự động nhận diện"])

    def analyze_pptx_and_integrate(self, slides_text: str, cap_hoc: str, integration_type: str) -> dict:
        """Phân tích các slide PowerPoint và đề xuất nội dung ghi chú diễn giả (Notes)."""
        tt02_info = self._get_tt02_framework_prompt(cap_hoc)
        ai_framework_info = """
* Khung Năng lực AI (QĐ 2422/QĐ-BGDĐT):
  - Nhận thức về AI: Nhận diện ứng dụng AI trong đời sống và học tập.
  - Ứng dụng AI: Dùng AI hỗ trợ gợi ý ý tưởng, tóm tắt, tra cứu thông tin, dịch thuật.
  - Tư duy phản biện & Đạo đức AI: Đánh giá độ tin cậy kết quả của AI, tôn trọng bản quyền và sử dụng có trách nhiệm.
"""

        if integration_type == "Năng lực số":
            focus_instruction = f"YÊU CẦU: CHỈ TÍCH HỢP NĂNG LỰC SỐ (TT 02/2025/TT-BGDĐT).\n{tt02_info}\nTất cả các mục có 'loai': 'Năng lực số'."
        elif integration_type == "Năng lực AI":
            focus_instruction = f"YÊU CẦU: CHỈ TÍCH HỢP NĂNG LỰC AI (QĐ 2422/QĐ-BGDĐT).\n{ai_framework_info}\nTất cả các mục có 'loai': 'Năng lực AI'."
        else:
            focus_instruction = f"""
YÊU CẦU BẮT BUỘC KHI CHỌN 'CẢ HAI':
Phải tích hợp cả NĂNG LỰC SỐ (TT 02/2025/TT-BGDĐT) và NĂNG LỰC AI (QĐ 2422/QĐ-BGDĐT).
Danh sách trả về phải có sự phân bổ cả 2 loại ('Năng lực số' và 'Năng lực AI') phù hợp với từng slide.
"""

        prompt = f"""
Bạn là chuyên gia sư phạm và chuyển đổi số trong giáo dục phổ thông Việt Nam.
Hãy đọc danh sách các Slide bài giảng PowerPoint dưới đây. Hãy phân tích nội dung từng slide và đề xuất các GHI CHÚ SƯ PHẠM (để đưa vào phần Slide Notes cho giáo viên) nhằm tích hợp Năng lực số / Năng lực AI vào hoạt động dạy học của slide đó.

Cấp học chỉ định: {cap_hoc}

{focus_instruction}

QUY TẮC QUAN TRỌNG:
1. Xác định chính xác `slide_number` (số nguyên) của slide cần tích hợp.
2. Nội dung `insert_content` là lời nhắc/hướng dẫn sư phạm ngắn gọn, thiết thực cho giáo viên (ví dụ: "GV hướng dẫn HS tra cứu dữ liệu số...", "GV nhắc HS dùng AI tạo gợi ý nhưng cần đối chiếu SGK...", "Tổ chức cho HS chia sẻ sản phẩm số qua link nhóm...").
3. Không cần tích hợp trên tất cả mọi slide, chỉ chọn những slide hoạt động trọng tâm, slide thảo luận hoặc bài tập.

Danh sách nội dung các Slide:
----------------------------------
{slides_text}
----------------------------------
"""

        try:
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt,
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=TichHopResult,
                    temperature=0.3
                )
            )
            return json.loads(response.text)
        except APIError as ae:
            raise RuntimeError(f"Lỗi kết nối Gemini API: {str(ae)}")
        except json.JSONDecodeError:
            raise RuntimeError("Gemini phản hồi sai cấu trúc JSON.")
        except Exception as e:
            raise RuntimeError(f"Đã xảy ra lỗi: {str(e)}")
```eof

```python:1_📝_Tích_Hợp_Kỹ_Năng_Số.py
import streamlit as st
from gemini_service import GeminiService
from word_processor import WordProcessor
from pptx_processor import PPTXProcessor

st.set_page_config(
    page_title="Tích hợp Năng lực số & AI vào KHBD / PowerPoint",
    page_icon="📝",
    layout="wide"
)

st.markdown("## 🤖 Tích hợp Năng lực số và AI tự động vào Giáo án / PowerPoint")
st.info("Hỗ trợ tự động chèn mục tiêu/hoạt động vào Giáo án Word (.docx) hoặc chèn hướng dẫn sư phạm vào Slide Notes của bài giảng PowerPoint (.pptx).")

# --- CẤU HÌNH HỆ THỐNG ---
with st.expander("⚙️ **CẤU HÌNH HỆ THỐNG:**", expanded=False):
    col_cfg1, col_cfg2, col_cfg3 = st.columns([2, 1, 1])
    
    with col_cfg1:
        if "gemini_api_key" in st.session_state and st.session_state["gemini_api_key"].strip() != "":
            api_key = st.session_state["gemini_api_key"]
            st.success("🔑 **Trạng thái API Key:** Đã nhận diện.")
        else:
            st.warning("⚠️ **Chưa tìm thấy API Key:** Vui lòng nhập tại Trang chủ.")
            st.stop()

    with col_cfg2:
        cap_hoc = st.selectbox(
            "**Chọn cấp học mục tiêu:**",
            ["Tự động nhận diện", "Tiểu học", "THCS", "THPT"]
        )

    with col_cfg3:
        integration_type = st.selectbox(
            "**Loại tích hợp:**",
            ["Cả hai", "Năng lực số", "Năng lực AI"]
        )

# --- MÀN HÌNH CHÍNH: 2 CỘT ---
col_left, col_right = st.columns([2, 1])

with col_left:
    with st.container(border=True):
        st.markdown("#### 📂 1. Tải lên tệp Giáo án (.docx) hoặc Bài giảng (.pptx)")
        uploaded_file = st.file_uploader(
            "**Chọn file (.docx hoặc .pptx):**", 
            type=["docx", "pptx"],
            help="Hệ thống hỗ trợ cả Word (.docx) và PowerPoint (.pptx)."
        )

        if uploaded_file is not None:
            file_bytes = uploaded_file.read()
            file_ext = uploaded_file.name.split('.')[-1].lower()
            st.success(f"✔️ Đã tải lên file: **{uploaded_file.name}**")
            st.session_state['original_filename'] = uploaded_file.name
            st.session_state['file_ext'] = file_ext
            
            if st.button("🚀 Bắt đầu tích hợp", type="primary", use_container_width=True):
                with st.spinner("🔄 Đang xử lý nội dung và gửi phân tích tới Gemini AI..."):
                    try:
                        ai_handler = GeminiService(api_key=api_key)
                        
                        if file_ext == "pptx":
                            # Xử lý file PowerPoint
                            slides_data = PPTXProcessor.extract_slides_text(file_bytes)
                            if not slides_data:
                                st.error("❌ Không tìm thấy văn bản trong file PowerPoint.")
                                st.stop()
                            
                            doc_text = PPTXProcessor.format_doc_text_for_ai(slides_data)
                            ai_result = ai_handler.analyze_pptx_and_integrate(doc_text, cap_hoc, integration_type)
                            st.session_state['ai_result'] = ai_result
                            
                            processed_file = PPTXProcessor.integrate_into_notes(file_bytes, ai_result)
                            st.session_state['processed_file'] = processed_file
                        else:
                            # Xử lý file Word
                            doc_text = WordProcessor.extract_text(file_bytes)
                            if not doc_text.strip():
                                st.error("❌ File Word trống hoặc không đọc được văn bản.")
                                st.stop()
                                
                            ai_result = ai_handler.analyze_and_integrate(doc_text, cap_hoc, integration_type)
                            st.session_state['ai_result'] = ai_result
                            
                            processed_file = WordProcessor.integrate_digital_capacity(file_bytes, ai_result, integration_type)
                            st.session_state['processed_file'] = processed_file
                        
                        st.success("🎉 Tích hợp thành công!")
                        
                    except Exception as e:
                        st.error(f"❌ Đã xảy ra lỗi trong quá trình xử lý: {str(e)}")

    # Hiển thị kết quả & Nút tải về
    if 'ai_result' in st.session_state and 'processed_file' in st.session_state:
        with st.container(border=True):
            st.markdown("#### 📋 2. Kết quả tích hợp")
            res = st.session_state['ai_result']
            sua_doi_list = res.get('sua_doi', [])
            file_ext = st.session_state.get('file_ext', 'docx')
            
            if not sua_doi_list:
                st.warning("AI không đề xuất vị trí tích hợp nào.")
            else:
                for idx, item in enumerate(sua_doi_list):
                    content = item.get('insert_content', 'Không có nội dung')
                    loai = item.get('loai', 'Năng lực số')
                    icon = "🧠" if loai == "Năng lực AI" else "💻"
                    color = "#D97706" if loai == "Năng lực AI" else "#0066CC"
                    
                    if file_ext == "pptx":
                        slide_num = item.get('slide_number', 'Chưa rõ')
                        title = f"{icon} Slide {slide_num} ({loai})"
                        with st.expander(title, expanded=True):
                            st.markdown(f"**Ghi chú bổ sung vào Slide Notes:** <span style='color:{color}; font-weight:bold;'>{content}</span>", unsafe_allow_html=True)
                    else:
                        anchor = item.get('anchor_text', 'Không rõ vị trí')
                        title = f"{icon} Vị trí {idx+1}: Sau \"{anchor}\" ({loai})"
                        with st.expander(title, expanded=True):
                            st.markdown(f"**Văn bản gốc:** `{anchor}`")
                            st.markdown(f"**Nội dung chèn:** <span style='color:{color}; font-weight:bold;'>{content}</span>", unsafe_allow_html=True)
            
            st.markdown("---")
            orig_name = st.session_state.get('original_filename', 'KHBD_TichHop')
            base_name = orig_name.rsplit('.', 1)[0]
            
            if file_ext == "pptx":
                download_filename = f"{base_name}_TichHop_Notes.pptx"
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                download_filename = f"{base_name}_Tichhop_So_AI.docx"
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            
            st.download_button(
                label=f"💾 TẢI XUỐNG TỆP ĐÃ TÍCH HỢP (.{file_ext.upper()})",
                data=st.session_state['processed_file'],
                file_name=download_filename,
                type="primary",
                mime=mime_type,
                use_container_width=True
            )

with col_right:
    st.markdown("### ℹ️ Hướng dẫn sử dụng")
    st.markdown("""
    - **Hỗ trợ 2 định dạng:**
      + **Word (.docx):** Chèn trực tiếp các dòng năng lực số & AI vào mục tiêu và tiến trình hoạt động.
      + **PowerPoint (.pptx):** Tự động thêm các hướng dẫn, lưu ý sư phạm vào phần **Slide Notes** của từng slide.
    - **Chuẩn khung tham chiếu:**
      + **Năng lực số:** Thông tư 02/2025/TT-BGDĐT.
      + **Năng lực AI:** Quyết định 2422/QĐ-BGDĐT.
    """)
```eof

### Tóm tắt các nâng cấp đã thực hiện:
1. **Thêm module `pptx_processor.py`**: Trích xuất nội dung văn bản từ các Slide PowerPoint và tự động ghi chú các hướng dẫn tích hợp vào `slide.notes_slide` mà không làm thay đổi giao diện slide.
2. **Cập nhật `gemini_service.py`**: Bổ sung hàm `analyze_pptx_and_integrate` với định dạng Schema hỗ trợ trường `slide_number`.
3. **Cập nhật giao diện `1_📝_Tích_Hợp_Kỹ_Năng_Số.py`**: Cho phép tải lên cả 2 định dạng `.docx` và `.pptx`, tự động nhận diện và tải về đúng định dạng tương ứng.

*(Lưu ý: Bạn hãy cài đặt thêm thư viện `python-pptx` nếu chưa có bằng lệnh: `pip install python-pptx`)*

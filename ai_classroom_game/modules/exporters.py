import pandas as pd
import io
from docx import Document
from pptx import Presentation

class DocumentExporter:
    @staticmethod
    def export_to_docx(topic, quiz_data):
        doc = Document()
        doc.add_heading(f"KỊCH BẢN HOẠT ĐỘNG: {topic.upper()}", level=1)
        
        doc.add_heading("1. Phần Câu Hỏi Trắc Nghiệm", level=2)
        for idx, q in enumerate(quiz_data.get("trac_nghiem", []), 1):
            doc.add_paragraph(f"Câu {idx}: {q['cau_hoi']}")
            for opt in q['options']:
                doc.add_paragraph(f"  [ ] {opt}")
            doc.add_paragraph(f"👉 Đáp án đúng: {q['dap_an']} - Giải thích: {q['giai_thich']}")
            
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    @staticmethod
    def export_to_pptx(topic, quiz_data):
        prs = Presentation()
        # Slide tiêu đề chính
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Trò Chơi Học Tập\nChủ đề: {topic}"
        
        # Tạo slide câu hỏi tự động
        for q in quiz_data.get("trac_nghiem", []):
            blank_layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(blank_layout)
            s.shapes.title.text = q['cau_hoi']
            body = s.placeholders[1]
            body.text = "\n".join(q['options'])
            
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    @staticmethod
    def export_to_platform_excel(quiz_data, platform="Quizizz"):
        """Xuất ma trận câu hỏi ra file Excel theo form chuẩn của Kahoot hoặc Quizizz"""
        buffer = io.BytesIO()
        rows = []
        
        # Duyệt qua bộ câu hỏi trắc nghiệm được sinh ra từ AI
        for idx, item in enumerate(quiz_data.get("trac_nghiem", [])):
            opts = item.get("options", ["", "", "", ""])
            # Đảm bảo đủ 4 phương án
            while len(opts) < 4:
                opts.append("")
                
            if platform == "Quizizz":
                # Định dạng cột chuẩn của Quizizz
                rows.append({
                    "Question Text": item.get("cau_hoi", ""),
                    "Question Type": "Multiple Choice",
                    "Option 1": opts[0],
                    "Option 2": opts[1],
                    "Option 3": opts[2],
                    "Option 4": opts[3],
                    "Correct Answer": opts.index(item.get("dap_an")) + 1 if item.get("dap_an") in opts else 1,
                    "Time in seconds": 30
                })
            else:
                # Định dạng cột chuẩn của Kahoot
                rows.append({
                    "Question - max 120 chars": item.get("cau_hoi", ""),
                    "Answer 1 - max 75 chars": opts[0],
                    "Answer 2 - max 75 chars": opts[1],
                    "Answer 3 - max 75 chars": opts[2],
                    "Answer 4 - max 75 chars": opts[3],
                    "Time limit (sec)": 30,
                    "Correct answer(s)": opts.index(item.get("dap_an")) + 1 if item.get("dap_an") in opts else 1
                })
                
        df = pd.DataFrame(rows)
        
        # Ghi đè xuất file Excel qua openpyxl
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Questions")
            
        return buffer.getvalue()

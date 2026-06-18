from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import requests

def create_powerpoint_file(slide_data_json, style_choice):
    """Xây dựng file .pptx, chuẩn hóa công thức và tự động chèn ảnh minh họa từ Internet"""
    prs = Presentation()
    
    # Định nghĩa cấu hình màu sắc theo cấp học
    if style_choice == "Tiểu học":
        title_color = RGBColor(2, 132, 199)
        text_color = RGBColor(30, 41, 59)
    elif style_choice == "THCS":
        title_color = RGBColor(30, 58, 138)
        text_color = RGBColor(15, 23, 42)
    else:
        title_color = RGBColor(17, 24, 39)
        text_color = RGBColor(55, 65, 81)

    slides_list = slide_data_json.get("slides", [])
    
    for s_idx, s_data in enumerate(slides_list):
        s_type = s_data.get("type", "content")
        
        if s_type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = s_data.get("title", "BÀI GIẢNG ĐIỆN TỬ")
            title.text_frame.paragraphs[0].font.color.rgb = title_color
            
            subtitle = slide.placeholders[1]
            sub_text = f"{s_data.get('subtitle', '')}\nGiáo viên: {s_data.get('teacher', 'Ngo Thanh Hung')}\nTrường học: {s_data.get('school', 'Trường THCS')}"
            subtitle.text = sub_text
            
        elif s_type == "quiz":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            title_shape.text = s_data.get("title", "Câu hỏi luyện tập")
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 38, 38)
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = s_data.get("question", "") + "\n"
            for opt in s_data.get("options", []):
                p = tf.add_paragraph()
                p.text = f"❑ {opt}"
                p.level = 1
                
        else: # Slide nội dung thông thường
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_shape = slide.shapes.title
            title_shape.text = s_data.get("title", "Nội dung bài học")
            title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
            
            # Thu hẹp chiều rộng khung chữ lại để nhường không gian bên phải cho ảnh
            body_shape = slide.placeholders[1]
            body_shape.width = Inches(5.5) 
            tf = body_shape.text_frame
            tf.text = ""
            
            raw_content = s_data.get("bullets", s_data.get("content", s_data.get("body", "")))
            
            if isinstance(raw_content, list):
                for idx, item in enumerate(raw_content):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.text = f"• {item}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = text_color
            else:
                tf.text = str(raw_content)
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = text_color

            # --- TỰ ĐỘNG LẤY VÀ CHÈN ẢNH TỪ INTERNET ---
            keyword = s_data.get("prompt_image", "")
            if keyword and len(keyword) > 2:
                try:
                    # Sử dụng dịch vụ ảnh gốc băm theo từ khóa (Ví dụ lấy từ source.unsplash)
                    # Thiết lập kích thước ảnh tải về gọn gàng 400x300
                    img_url = f"https://images.unsplash.com/photo-1546410531-bb4caa6b424d?w=400" # Ảnh mặc định giáo dục
                    
                    # Tạo link tìm kiếm ảnh động theo từ khóa của AI (đã lược bớt ký tự đặc biệt)
                    clean_kw = "".join(c for c in keyword if c.isalnum() or c==' ').strip().replace(" ", ",")
                    dynamic_url = f"https://source.unsplash.com/featured/400x300/?{clean_kw}"
                    
                    response = requests.get(dynamic_url, timeout=5)
                    if response.status_code == 200:
                        image_stream = io.BytesIO(response.content)
                        # Chèn ảnh vào tọa độ: Cách lề trái 6.2 Inches, Cách lề trên 2.0 Inches
                        slide.shapes.add_picture(image_stream, Inches(6.2), Inches(2.0), width=Inches(3.4))
                except Exception:
                    pass # Nếu lỗi mạng hoặc không tìm thấy ảnh thì bỏ qua để tránh sập app

        # Tích hợp Teacher Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = s_data.get("notes", "Giáo viên giảng dạy theo tiến trình slide.")

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
